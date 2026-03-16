"""Chart slide catalog for template2.pptx.

Each chart type maps to a template slide index and describes what data
fields the LLM should generate.  The fill functions add text boxes
programmatically on top of the decorative graphics.
"""

from typing import Any

from pptx.util import Emu
from pptx.dml.color import RGBColor

# ──────────────────────────────────────────────────────────
# Reuse constants from pptx_builder (imported at runtime)
# ──────────────────────────────────────────────────────────
FONT_NAME = "方正银联黑简体"
TEXT_COLOR = RGBColor(65, 65, 65)
SUBTITLE_COLOR = RGBColor(130, 130, 130)
ACCENT_COLOR = RGBColor(244, 58, 62)

# ──────────────────────────────────────────────────────────
# Chart catalog — slide index + description
# ──────────────────────────────────────────────────────────
CHART_CATALOG = {
    "arrow_flow": {
        "template_index": 5,
        "description": "5-step arrow flow chart with labels and descriptions",
    },
    "timeline": {
        "template_index": 7,
        "description": "Horizontal timeline with 4 milestone nodes",
    },
    "quadrant": {
        "template_index": 8,
        "description": "4-quadrant analysis with titles and content",
    },
    "card_4col": {
        "template_index": 9,
        "description": "4-column card layout with banner, titles, and content",
    },
    "cycle_4": {
        "template_index": 11,
        "description": "4-node cycle diagram with labels and content",
    },
    "compare_list": {
        "template_index": 10,
        "description": "3-column comparison list with headers and content",
    },
    "pyramid": {
        "template_index": 12,
        "description": "3-level pyramid hierarchy with labels and content",
    },
}

# ──────────────────────────────────────────────────────────
# Prompt section for LLM
# ──────────────────────────────────────────────────────────
CHART_PROMPT_SECTION = """
CHART LAYOUTS (use when content fits a visual pattern):

8. "arrow_flow" — 5-step arrow flow diagram.
   Fields: "title", "steps" (array of 5 objects with "label" and "content")
   Use: Process flows, value chains, sequential steps.

9. "timeline" — Horizontal timeline with 4 milestones.
   Fields: "title", "milestones" (array of 4 strings)
   Use: Chronological events, development history, project phases.

10. "quadrant" — 4-quadrant analysis chart.
    Fields: "title", "quadrants" (array of 4 objects with "label" and "content")
    Use: SWOT analysis, 2×2 matrix, categorization.

11. "card_4col" — 4 column cards with shared banner.
    Fields: "title", "banner" (subtitle text), "cards" (array of 4 objects with "label" and "content")
    Use: Feature comparison, service overview, key metrics.

12. "cycle_4" — 4-node cycle diagram.
    Fields: "title", "nodes" (array of 4 objects with "label" and "content")
    Use: Iterative processes, feedback loops, cyclical workflows.

13. "compare_list" — 3-column comparison table.
    Fields: "title", "columns" (array of 3 objects with "header" and "items" array of strings)
    Use: Feature comparison, before/after, multi-option evaluation.

14. "pyramid" — 3-level pyramid hierarchy.
    Fields: "title", "levels" (array of 3 objects with "label" and "content", top to bottom)
    Use: Hierarchy, importance ranking, layered architecture.
"""


# ──────────────────────────────────────────────────────────
# Fill functions — each adds text programmatically
# ──────────────────────────────────────────────────────────

def _set_simple(shape, text, font_size=14, font_color=TEXT_COLOR, bold=False):
    """Set shape text with simple formatting."""
    tf = shape.text_frame
    tf.clear()
    from pptx.util import Pt
    p = tf.paragraphs[0]
    p.text = str(text)
    run = p.runs[0] if p.runs else p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.name = FONT_NAME
    run.font.color.rgb = font_color
    run.font.bold = bold


def _set_body(shape, text, font_size=12, font_color=TEXT_COLOR, line_spacing=1.2):
    """Set shape text with line spacing for body content."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    from pptx.util import Pt
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(2)
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.name = FONT_NAME
            run.font.color.rgb = font_color


def _replace_header(slide, title):
    """Replace the '图表示例' header with slide title."""
    if not title:
        return
    for shape in slide.shapes:
        if (shape.has_text_frame
                and "Placeholder" in shape.name
                and shape.text_frame.text.strip() in ("图表示例", "图表内容")):
            _set_simple(shape, title, font_size=24, font_color=SUBTITLE_COLOR)
            return


# ─── arrow_flow (Slide 6) ────────────────────────────────
def fill_arrow_flow(slide: Any, data: dict) -> None:
    """5-step arrow flow: replace labels (TextBox 81) and descriptions (TextBox 64)."""
    _replace_header(slide, data.get("title"))

    steps = data.get("steps", [])
    if not steps:
        return

    # Labels: 5 × TextBox 81 (sorted by position)
    labels = sorted(
        [s for s in slide.shapes if s.has_text_frame and s.name == "TextBox 81"],
        key=lambda s: (s.top, s.left),
    )
    # Descriptions: TextBox 64 (sorted) + 文本框 69 (center)
    descs = sorted(
        [s for s in slide.shapes if s.has_text_frame and s.name in ("TextBox 64", "文本框 69")],
        key=lambda s: (s.top, s.left),
    )

    for i, step in enumerate(steps[:5]):
        if isinstance(step, dict):
            label = step.get("label", "")
            content = step.get("content", "")
        else:
            label = str(step)
            content = ""

        if i < len(labels) and label:
            _set_simple(labels[i], label, font_size=12, font_color=ACCENT_COLOR, bold=True)
        if i < len(descs) and content:
            _set_body(descs[i], content, font_size=10)


# ─── timeline (Slide 13) ─────────────────────────────────
def fill_timeline(slide: Any, data: dict) -> None:
    """4-milestone timeline: replace 文本框 59/60/61/62."""
    _replace_header(slide, data.get("title"))

    milestones = data.get("milestones", [])
    if not milestones:
        return

    targets = sorted(
        [s for s in slide.shapes
         if s.has_text_frame and s.name in ("文本框 59", "文本框 60", "文本框 61", "文本框 62")],
        key=lambda s: s.left,  # left-to-right order
    )

    for i, ms in enumerate(milestones[:4]):
        if i < len(targets):
            text = ms if isinstance(ms, str) else ms.get("label", str(ms))
            _set_simple(targets[i], text, font_size=12, font_color=TEXT_COLOR, bold=True)


# ─── quadrant (Slide 19) ─────────────────────────────────
def fill_quadrant(slide: Any, data: dict) -> None:
    """4-quadrant: replace 矩形 58/60/62/64 (labels) + 文本框 36/37/38/65 (content)."""
    _replace_header(slide, data.get("title"))

    quads = data.get("quadrants", [])
    if not quads:
        return

    label_shapes = sorted(
        [s for s in slide.shapes
         if s.has_text_frame and s.name in ("矩形 58", "矩形 60", "矩形 62", "矩形 64")],
        key=lambda s: (s.top, s.left),
    )
    content_shapes = sorted(
        [s for s in slide.shapes
         if s.has_text_frame and s.name in ("文本框 36", "文本框 37", "文本框 38", "文本框 65")],
        key=lambda s: (s.top, s.left),
    )

    for i, q in enumerate(quads[:4]):
        if isinstance(q, dict):
            label = q.get("label", "")
            content = q.get("content", "")
        else:
            label = str(q)
            content = ""

        if i < len(label_shapes) and label:
            _set_simple(label_shapes[i], label, font_size=12, font_color=RGBColor(255, 255, 255), bold=True)
        if i < len(content_shapes) and content:
            _set_body(content_shapes[i], content, font_size=10)


# ─── card_4col (Slide 23) ────────────────────────────────
def fill_card_4col(slide: Any, data: dict) -> None:
    """4-column cards: replace 矩形 36 (banner), 矩形 38/39/40/41 (titles), Rectangle 20 (content)."""
    _replace_header(slide, data.get("title"))

    # Banner
    for s in slide.shapes:
        if s.has_text_frame and s.name == "矩形 36":
            _set_simple(s, data.get("banner", ""), font_size=14, font_color=RGBColor(255, 255, 255))
            break

    cards = data.get("cards", [])
    if not cards:
        return

    # Card title shapes (矩形 38/39/40/41)
    title_shapes = sorted(
        [s for s in slide.shapes
         if s.has_text_frame and s.name in ("矩形 38", "矩形 39", "矩形 40", "矩形 41")],
        key=lambda s: (s.top, s.left),
    )
    # Card content shapes (Rectangle 20)
    content_shapes = sorted(
        [s for s in slide.shapes
         if s.has_text_frame and s.name == "Rectangle 20"],
        key=lambda s: (s.top, s.left),
    )

    for i, card in enumerate(cards[:4]):
        if isinstance(card, dict):
            label = card.get("label", "")
            content = card.get("content", "")
        else:
            label = str(card)
            content = ""

        if i < len(title_shapes) and label:
            _set_simple(title_shapes[i], label, font_size=12, font_color=RGBColor(255, 255, 255), bold=True)
        if i < len(content_shapes) and content:
            _set_body(content_shapes[i], content, font_size=10)


# ─── numbered_4 (Slide 47) ───────────────────────────────
def fill_numbered_4(slide: Any, data: dict) -> None:
    """4 numbered steps: TextBox 14-17 (nums), TextBox 22-29 (title+content pairs)."""
    _replace_header(slide, data.get("title"))

    steps = data.get("steps", [])
    if not steps:
        return

    # Number shapes (TextBox 14/15/16/17)
    nums = sorted(
        [s for s in slide.shapes
         if s.has_text_frame and s.name in ("TextBox 14", "TextBox 15", "TextBox 16", "TextBox 17")],
        key=lambda s: s.left,
    )
    # Title shapes (TextBox 22/24/26/28)
    titles = sorted(
        [s for s in slide.shapes
         if s.has_text_frame and s.name in ("TextBox 22", "TextBox 24", "TextBox 26", "TextBox 28")],
        key=lambda s: s.left,
    )
    # Content shapes (TextBox 23/25/27/29)
    contents = sorted(
        [s for s in slide.shapes
         if s.has_text_frame and s.name in ("TextBox 23", "TextBox 25", "TextBox 27", "TextBox 29")],
        key=lambda s: s.left,
    )

    for i, step in enumerate(steps[:4]):
        if isinstance(step, dict):
            label = step.get("label", "")
            content = step.get("content", "")
        else:
            label = str(step)
            content = ""

        if i < len(nums):
            _set_simple(nums[i], f"{i+1:02d}", font_size=28, font_color=ACCENT_COLOR, bold=True)
        if i < len(titles) and label:
            _set_simple(titles[i], label, font_size=12, font_color=TEXT_COLOR, bold=True)
        if i < len(contents) and content:
            _set_body(contents[i], content, font_size=10)


# ─── cycle_4 (Slide 40) ──────────────────────────────────
def fill_cycle_4(slide: Any, data: dict) -> None:
    """4-node cycle: Shape 376 (labels) + Shape 357 (content)."""
    _replace_header(slide, data.get("title"))

    nodes = data.get("nodes", [])
    if not nodes:
        return

    labels = sorted(
        [s for s in slide.shapes if s.has_text_frame and s.name == "Shape 376"],
        key=lambda s: s.left,
    )
    contents = sorted(
        [s for s in slide.shapes if s.has_text_frame and s.name == "Shape 357"],
        key=lambda s: s.left,
    )

    for i, node in enumerate(nodes[:4]):
        if isinstance(node, dict):
            label = node.get("label", "")
            content = node.get("content", "")
        else:
            label = str(node)
            content = ""

        if i < len(labels) and label:
            _set_simple(labels[i], label, font_size=10, font_color=RGBColor(255, 255, 255), bold=True)
        if i < len(contents) and content:
            _set_body(contents[i], content, font_size=9)


# ─── compare_list (Slide 30) ─────────────────────────────
def fill_compare_list(slide: Any, data: dict) -> None:
    """3-column comparison: header shapes + content shapes."""
    _replace_header(slide, data.get("title"))

    columns = data.get("columns", [])
    if not columns:
        return

    # Header shapes: 文本框 28, 矩形 31, 矩形 33 (left to right)
    headers = sorted(
        [s for s in slide.shapes
         if s.has_text_frame and s.name in ("文本框 28", "矩形 31", "矩形 33")],
        key=lambda s: s.left,
    )
    # Content shapes: 矩形 53, 矩形 30, 矩形 54 (left to right)
    bodies = sorted(
        [s for s in slide.shapes
         if s.has_text_frame and s.name in ("矩形 53", "矩形 30", "矩形 54")],
        key=lambda s: s.left,
    )

    for i, col in enumerate(columns[:3]):
        if isinstance(col, dict):
            header = col.get("header", "")
            items = col.get("items", [])
            content = "\n".join(items) if isinstance(items, list) else str(items)
        else:
            header = str(col)
            content = ""

        if i < len(headers) and header:
            _set_simple(headers[i], header, font_size=12, font_color=TEXT_COLOR, bold=True)
        if i < len(bodies) and content:
            _set_body(bodies[i], content, font_size=10)


# ─── pyramid (Slide 44) ──────────────────────────────────
def fill_pyramid(slide: Any, data: dict) -> None:
    """3-level pyramid: group children Rectangle 39/40/41 + Rectangle 14/15."""
    _replace_header(slide, data.get("title"))

    levels = data.get("levels", [])
    if not levels:
        return

    # Pyramid labels: Rectangle 41 (top standalone), Rectangle 40 (in Group 31),
    # Rectangle 39 (in Group 33)
    # Since these are inside groups, we add text boxes programmatically instead
    positions = [
        # (x, y, w, h) for top, middle, bottom labels
        (Emu(int(11.2 * 914400 / 2.54)), Emu(int(7.3 * 914400 / 2.54)),
         Emu(int(3.0 * 914400 / 2.54)), Emu(int(0.9 * 914400 / 2.54))),
        (Emu(int(9.5 * 914400 / 2.54)), Emu(int(9.5 * 914400 / 2.54)),
         Emu(int(3.0 * 914400 / 2.54)), Emu(int(0.9 * 914400 / 2.54))),
        (Emu(int(7.8 * 914400 / 2.54)), Emu(int(11.7 * 914400 / 2.54)),
         Emu(int(3.0 * 914400 / 2.54)), Emu(int(0.9 * 914400 / 2.54))),
    ]
    content_positions = [
        (Emu(int(14.5 * 914400 / 2.54)), Emu(int(7.0 * 914400 / 2.54)),
         Emu(int(8.0 * 914400 / 2.54)), Emu(int(2.0 * 914400 / 2.54))),
        (Emu(int(14.5 * 914400 / 2.54)), Emu(int(9.2 * 914400 / 2.54)),
         Emu(int(8.0 * 914400 / 2.54)), Emu(int(2.0 * 914400 / 2.54))),
        (Emu(int(14.5 * 914400 / 2.54)), Emu(int(11.4 * 914400 / 2.54)),
         Emu(int(8.0 * 914400 / 2.54)), Emu(int(2.0 * 914400 / 2.54))),
    ]

    # Try to fill the standalone Rectangle 41 (top label)
    for s in slide.shapes:
        if s.has_text_frame and s.name == "Rectangle 41" and levels:
            level = levels[0]
            label = level.get("label", "") if isinstance(level, dict) else str(level)
            _set_simple(s, label, font_size=12, font_color=RGBColor(255, 255, 255), bold=True)

    # Fill group labels (Rectangle 39, 40) inside groups
    label_names = {"Rectangle 39": 2, "Rectangle 40": 1}  # name -> level index
    for shape in slide.shapes:
        if shape.shape_type == 6 and hasattr(shape, 'shapes'):
            for child in shape.shapes:
                if child.has_text_frame and child.name in label_names:
                    idx = label_names[child.name]
                    if idx < len(levels):
                        level = levels[idx]
                        label = level.get("label", "") if isinstance(level, dict) else str(level)
                        _set_simple(child, label, font_size=12, font_color=RGBColor(255, 255, 255), bold=True)

    # Add content text boxes
    for i, level in enumerate(levels[:3]):
        if isinstance(level, dict) and level.get("content"):
            if i < len(content_positions):
                x, y, w, h = content_positions[i]
                txBox = slide.shapes.add_textbox(x, y, w, h)
                _set_body(txBox, level["content"], font_size=11, font_color=TEXT_COLOR)


# ──────────────────────────────────────────────────────────
# Registry
# ──────────────────────────────────────────────────────────
CHART_FILL_FUNCTIONS = {
    "arrow_flow": fill_arrow_flow,
    "timeline": fill_timeline,
    "quadrant": fill_quadrant,
    "card_4col": fill_card_4col,
    "cycle_4": fill_cycle_4,
    "compare_list": fill_compare_list,
    "pyramid": fill_pyramid,
}
