"""Build PPTX files by cloning template slides and replacing text content."""

import copy
import io
import os
from typing import Any

from loguru import logger
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from open_notebook.utils.chart_catalog import (
    CHART_CATALOG, CHART_FILL_FUNCTIONS,
)

# Path to our custom template (cleaned 银联品牌模板)
TEMPLATE_PATH = os.path.join(
    os.path.dirname(__file__), "..", "..", "ppttemplate", "template1.pptx"
)

# ──────────────────────────────────────────────────────────
# 银联模板格式规范（来源：templateyinlian2.pptx）
# ──────────────────────────────────────────────────────────
FONT_NAME = "方正银联黑简体"

# ─── 颜色 ───
TEXT_COLOR = RGBColor(65, 65, 65)          # 正文 / 标题 灰色80%
SUBTITLE_COLOR = RGBColor(130, 130, 130)   # 目录标题 / 页面顶端标题 / 序号
ACCENT_COLOR = RGBColor(244, 58, 62)       # 强调文字 红色
ENDING_TITLE_COLOR = RGBColor(130, 130, 130)

# ─── 字号 ───
COVER_TITLE_SIZE = Pt(48)    # 封面标题
TOC_TITLE_SIZE = Pt(54)      # 目录标题（加粗）
PAGE_HEADER_SIZE = Pt(36)    # 页面顶端标题
HEADING1_SIZE = Pt(24)       # 一级标题（加粗，■ 方形符号）
HEADING2_SIZE = Pt(18)       # 二级标题（□ 空心方形符号）
BODY_FONT_SIZE = Pt(20)      # 正文
BADGE_FONT_SIZE = Pt(15)     # badge 标签
CONTACT_SIZE = Pt(20)        # 联系人信息
ENDING_TITLE_SIZE = Pt(66)   # 结束词

# ─── 行距 / 段距 ───
LINE_SPACING = 1.5           # 1.5 倍行距
SPACE_AFTER_PT = Pt(0)       # 段后 0.3 磅（Pt 取整为 0）

# Template slide catalog (0-indexed in template2.pptx)
# Compact template with only the 14 used slides.

SLIDE_CATALOG = {
    "cover": {
        "template_index": 0,
        "description": "Cover page with title and subtitle",
    },
    "content": {
        "template_index": 1,
        "description": "Simple content slide for body text and lists (primary)",
    },
    "content_badges": {
        "template_index": 4,
        "description": "Content slide with badge labels and text body",
    },
    "content_alt": {
        "template_index": 3,
        "description": "Content slide with four title-content blocks",
    },
    "two_blocks": {
        "template_index": 6,
        "description": "Two content rows with titles",
    },
    "ending": {
        "template_index": 2,
        "description": "Thank you / ending page",
    },
}


def _clone_slide(prs: Any, source_slide: Any) -> Any:
    """Clone a slide within the presentation, preserving all shapes and images."""
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Remove default shapes from new slide
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("sp", "pic", "grpSp", "cxnSp", "graphicFrame"):
            sp_tree.remove(child)

    # Copy all shapes from source
    src_tree = source_slide.shapes._spTree
    for child in src_tree:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("sp", "pic", "grpSp", "cxnSp", "graphicFrame"):
            sp_tree.append(copy.deepcopy(child))

    # Copy image relationships
    for rel in source_slide.part.rels.values():
        if "image" in rel.reltype.lower():
            new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)

    return new_slide


def _get_text_shapes(slide: Any) -> list[dict]:
    """Get all text-bearing shapes with their metadata, sorted by position."""
    shapes = []
    for i, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        is_ph = shape.is_placeholder
        ph_idx = shape.placeholder_format.idx if is_ph else -1
        shapes.append({
            "index": i,
            "shape": shape,
            "name": shape.name,
            "ph_idx": ph_idx,
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
            "text": shape.text_frame.text,
        })
    return shapes


def _replace_header_group(slide: Any, title: str) -> None:
    """Replace the page header ('图表示例') with the slide title.

    After user template edits, the header is now a standalone
    'Text Placeholder 1' shape at the top of each chart slide.
    Falls back to the old 组 31 group structure if present.
    """
    if not title:
        return

    # Strategy 1: standalone Text Placeholder 1 (new template layout)
    for shape in slide.shapes:
        if (shape.has_text_frame
                and "Placeholder" in shape.name
                and shape.text_frame.text.strip() in ("图表示例", "图表内容")):
            _set_simple_text(shape, title, font_size=24,
                             font_color=SUBTITLE_COLOR)
            logger.debug(f"Replaced standalone header with title: {title}")
            return

    # Strategy 2: 组 31 group (legacy template layout)
    for shape in slide.shapes:
        if shape.shape_type == 6 and "31" in shape.name:
            try:
                for child in shape.shapes:
                    name = child.name
                    if child.has_text_frame:
                        if "Placeholder" in name:
                            _set_simple_text(child, title, font_size=36,
                                             font_color=SUBTITLE_COLOR)
                        elif "矩形" in name:
                            child.text_frame.clear()
                    if "直接连接符" in name or "Connector" in name:
                        try:
                            child.width = 0
                            child.height = 0
                        except Exception:
                            pass
                logger.debug(f"Replaced header group with title: {title}")
                return
            except Exception as e:
                logger.warning(f"Failed to replace header group: {e}")


def _set_shape_text(shape: Any, text: str, font_size: int = 20,
                    font_name: str = FONT_NAME,
                    font_color: RGBColor = TEXT_COLOR,
                    line_spacing: float = LINE_SPACING) -> None:
    """Set text in a shape, handling bullet points. Uses 银联 style defaults."""
    tf = shape.text_frame
    tf.clear()
    lines = text.strip().split("\n")
    first = True
    for line in lines:
        line = line.strip()
        if not line:
            continue
        clean = line.lstrip("-*•■□").strip()
        if not clean:
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = clean
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = font_color
        p.line_spacing = line_spacing
    logger.debug(f"Set shape text ({shape.name}): {tf.text[:50]}")


def _set_simple_text(shape: Any, text: str, font_size: int = 20,
                     font_name: str = FONT_NAME,
                     font_color: RGBColor = TEXT_COLOR,
                     bold: bool = False,
                     line_spacing: float | None = None) -> None:
    """Set simple text (no bullet processing) with 银联 style."""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = font_color
    p.font.bold = bold
    if line_spacing is not None:
        p.line_spacing = line_spacing


def _fill_cover(slide: Any, data: dict) -> None:
    """Fill cover slide: 标题 1(ph0) + 内容占位符 8(ph1)."""
    shapes = _get_text_shapes(slide)
    for s in shapes:
        if s["ph_idx"] == 0 and data.get("title"):
            _set_simple_text(s["shape"], data["title"],
                             font_size=48, font_color=TEXT_COLOR)
        elif s["ph_idx"] == 1 and data.get("subtitle"):
            _set_simple_text(s["shape"], data["subtitle"],
                             font_size=20, font_color=TEXT_COLOR)


def _fill_content_simple(slide: Any, data: dict) -> None:
    """Fill simple content slide (Slide 2): add title + body text boxes.

    Slide 2 is a blank branded page. We programmatically add:
      - A title text box at top-left (24pt, bold)
      - A body text box below for content (20pt, 1.5x line spacing)
    """
    title = data.get("title", "")
    content = data.get("content", "")

    if title:
        txBox = slide.shapes.add_textbox(
            Emu(647564), Emu(368660),   # same position as 文本框 9 header
            Emu(7800000), Emu(600000),
        )
        _set_simple_text(txBox, title, font_size=24,
                         font_color=SUBTITLE_COLOR, bold=True)

    if content:
        txBox = slide.shapes.add_textbox(
            Emu(647564), Emu(1100000),
            Emu(7800000), Emu(5000000),
        )
        _set_shape_text(txBox, content, font_size=20, font_color=TEXT_COLOR)


def _fill_content_badges(slide: Any, data: dict) -> None:
    """Fill content slide (Slide 5): 圆角矩形 badges + programmatic content box.

    Template small text boxes removed. We fill badge shapes and add a
    content text box at the right side.
    """
    shapes = _get_text_shapes(slide)
    logger.debug(f"_fill_content_badges: {len(shapes)} text shapes")

    # Replace header group (附录|图表示例 → title)
    _replace_header_group(slide, data.get("title", ""))

    # Fill badge shapes (圆角矩形)
    badges = sorted(
        [s for s in shapes if "圆角矩形" in s["name"]],
        key=lambda s: (s["top"], s["left"]),
    )
    if badges and data.get("title"):
        _set_simple_text(badges[0]["shape"], data["title"],
                         font_size=12, font_color=RGBColor(255, 255, 255),
                         bold=True)
    if len(badges) >= 2 and data.get("badge"):
        _set_simple_text(badges[1]["shape"], data["badge"],
                         font_size=12, font_color=RGBColor(255, 255, 255),
                         bold=True)

    # Add content text box on the right side
    if data.get("content"):
        txBox = slide.shapes.add_textbox(
            Emu(4800000), Emu(5200000),  # right of center
            Emu(7200000), Emu(4500000),
        )
        _set_shape_text(txBox, data["content"],
                        font_size=14, font_color=TEXT_COLOR)


def _fill_content_alt(slide: Any, data: dict) -> None:
    """Fill content_alt slide (Slide 4): decorative graphics + programmatic text.

    Template small text boxes removed. Add title + content text boxes.
    """
    # Replace header group (附录|图表示例 → title)
    _replace_header_group(slide, data.get("title", ""))

    # Add content text box
    if data.get("content"):
        txBox = slide.shapes.add_textbox(
            Emu(4200000), Emu(5800000),
            Emu(7800000), Emu(4500000),
        )
        _set_shape_text(txBox, data["content"],
                        font_size=14, font_color=TEXT_COLOR)


def _fill_two_blocks(slide: Any, data: dict) -> None:
    """Fill two-block slide (Slide 11): decorative graphics + programmatic text.

    Template small text boxes removed. Add left/right title+content text boxes.
    """
    # Replace header group (附录|图表示例 → title)
    _replace_header_group(slide, data.get("title", ""))

    # Left block
    if data.get("left_title"):
        txBox = slide.shapes.add_textbox(
            Emu(1600000), Emu(5500000),
            Emu(4800000), Emu(600000),
        )
        _set_simple_text(txBox, data["left_title"],
                         font_size=16, font_color=TEXT_COLOR, bold=True)
    if data.get("left"):
        txBox = slide.shapes.add_textbox(
            Emu(1600000), Emu(6200000),
            Emu(4800000), Emu(3600000),
        )
        _set_shape_text(txBox, data["left"],
                        font_size=14, font_color=TEXT_COLOR)

    # Right block
    if data.get("right_title"):
        txBox = slide.shapes.add_textbox(
            Emu(7000000), Emu(5500000),
            Emu(4800000), Emu(600000),
        )
        _set_simple_text(txBox, data["right_title"],
                         font_size=16, font_color=TEXT_COLOR, bold=True)
    if data.get("right"):
        txBox = slide.shapes.add_textbox(
            Emu(7000000), Emu(6200000),
            Emu(4800000), Emu(3600000),
        )
        _set_shape_text(txBox, data["right"],
                        font_size=14, font_color=TEXT_COLOR)


def _fill_n_points(slide: Any, data: dict, n: int) -> None:
    """Fill N-point slide (Slide 10): decorative graphics + programmatic text.

    Template small text boxes removed. Add N point text boxes.
    """
    # Replace header group (附录|图表示例 → title)
    _replace_header_group(slide, data.get("title", ""))

    points = data.get("points", [])
    if not points:
        return

    # Layout N points in a row below the chart graphic
    num = min(n, len(points))
    total_width = 8400000  # ~21cm usable width
    box_width = total_width // num
    start_x = 1600000
    start_y = 5500000

    for i in range(num):
        point = points[i]
        x = start_x + i * box_width

        if isinstance(point, dict):
            label = point.get("label", str(i + 1))
            content = point.get("content", "")
        else:
            label = str(i + 1)
            content = str(point)

        # Label (bold)
        txBox = slide.shapes.add_textbox(
            Emu(x), Emu(start_y),
            Emu(box_width - 100000), Emu(400000),
        )
        _set_simple_text(txBox, label, font_size=14,
                         font_color=ACCENT_COLOR, bold=True)

        # Content
        if content:
            txBox = slide.shapes.add_textbox(
                Emu(x), Emu(start_y + 450000),
                Emu(box_width - 100000), Emu(3500000),
            )
            _set_shape_text(txBox, content,
                            font_size=12, font_color=TEXT_COLOR)


def _fill_ending(slide: Any, data: dict) -> None:
    """Fill ending slide (Slide 3): 文本框 1("谢  谢！")."""
    shapes = _get_text_shapes(slide)
    for s in shapes:
        if s["name"] == "文本框 1":
            title = data.get("title", "谢  谢！")
            _set_simple_text(s["shape"], title, font_size=66,
                             font_color=ENDING_TITLE_COLOR)


FILL_FUNCTIONS = {
    "cover": _fill_cover,
    "content": _fill_content_simple,
    "content_badges": _fill_content_badges,
    "content_alt": _fill_content_alt,
    "two_blocks": _fill_two_blocks,
    "ending": _fill_ending,
}

# Merge chart catalog entries
SLIDE_CATALOG.update(CHART_CATALOG)
FILL_FUNCTIONS.update(CHART_FILL_FUNCTIONS)


def build_pptx(slides: list[dict]) -> bytes:
    """
    Build a PPTX file by cloning template slides and replacing text.

    Args:
        slides: List of dicts from LLM, each with 'layout' and content fields.

    Returns:
        bytes of the generated .pptx file.
    """
    template_path = os.path.abspath(TEMPLATE_PATH)
    if not os.path.exists(template_path):
        logger.error(f"Template not found at {template_path}")
        raise FileNotFoundError(f"PPT template not found: {template_path}")

    # Load template to use as source for cloning
    source_prs = Presentation(template_path)
    source_slides = list(source_prs.slides)
    logger.info(f"Loaded template: {len(source_slides)} source slides")

    # Create output presentation from the same template, then clear it
    prs = Presentation(template_path)

    nsmap = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }
    sld_id_lst = prs.element.find(".//p:sldIdLst", nsmap)
    if sld_id_lst is not None:
        for sld_id in list(sld_id_lst):
            rId = sld_id.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            sld_id_lst.remove(sld_id)
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except Exception:
                    pass
    logger.info("Cleared output presentation")

    for i, slide_data in enumerate(slides):
        layout_name = slide_data.get("layout", "content")

        # Handle removed layout types gracefully
        if layout_name == "two_rows":
            layout_name = "two_blocks"

        catalog_entry = SLIDE_CATALOG.get(layout_name)

        if not catalog_entry:
            logger.warning(f"Unknown layout '{layout_name}' at slide {i}, defaulting to 'content'")
            layout_name = "content"
            catalog_entry = SLIDE_CATALOG["content"]

        template_idx = catalog_entry["template_index"]
        if template_idx >= len(source_slides):
            logger.warning(f"Template index {template_idx} out of range, using 'content'")
            template_idx = SLIDE_CATALOG["content"]["template_index"]
            layout_name = "content"

        # Clone the template slide
        source_slide = source_slides[template_idx]
        new_slide = _clone_slide(prs, source_slide)

        # Fill content
        fill_fn = FILL_FUNCTIONS.get(layout_name)
        if fill_fn:
            try:
                fill_fn(new_slide, slide_data)
                logger.debug(f"Filled slide {i} ({layout_name})")
            except Exception as e:
                logger.error(f"Error filling slide {i} ({layout_name}): {e}", exc_info=True)

    logger.info(f"Built PPTX with {len(prs.slides)} slides")

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
